import os
import io
import json
from datetime import datetime
from typing import Optional, List
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, black, grey
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
from reportlab.pdfgen import canvas
from reportlab.platypus import PageTemplate, Frame, BaseDocTemplate

load_dotenv()

# --- Configuration ---
try:
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
    if not GOOGLE_API_KEY:
        raise ValueError("CRITICAL: GOOGLE_API_KEY environment variable not set.")
    genai.configure(api_key=GOOGLE_API_KEY)
except Exception as e:
    print(f"Error configuring Gemini API: {e}")

app = FastAPI(
    title="Personalized Study Material Generator",
    description="Generate targeted study materials based on student weak spots from PowerPoint presentations",
    version="3.0.0"
)

# --- CORS Middleware ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Study Material Types ---
MATERIAL_TYPES = {
    "comprehensive": {
        "name": "Comprehensive Study Notes",
        "description": "Detailed explanations covering all aspects",
        "prompt_modifier": "Create extremely detailed and comprehensive notes with thorough explanations, multiple examples, and deep conceptual understanding."
    },
    "summarized": {
        "name": "Summarized Key Points",
        "description": "Concise summary focusing on essential concepts",
        "prompt_modifier": "Create a concise summary that captures only the most essential information and key takeaways. Be brief but complete."
    },
    "pointwise": {
        "name": "Point-wise Structure",
        "description": "Bullet points for easy memorization",
        "prompt_modifier": "Structure everything in clear bullet points and numbered lists. Make it easy to scan and memorize. Use short, impactful statements."
    },
    "prerequisite": {
        "name": "Prerequisite-Focused",
        "description": "Emphasizes foundational concepts",
        "prompt_modifier": "Focus heavily on prerequisite knowledge and foundational concepts. Build from basics to advanced. Include background information that students might be missing."
    },
    "problem_solving": {
        "name": "Problem-Solving Oriented",
        "description": "Focus on practical applications and exercises",
        "prompt_modifier": "Emphasize problem-solving techniques, worked examples, practice problems, and step-by-step solutions. Include common mistakes to avoid."
    },
    "visual_learning": {
        "name": "Visual Learning Guide",
        "description": "Structured for visual learners with diagrams descriptions",
        "prompt_modifier": "Structure content for visual learners. Describe concepts that would benefit from diagrams, flowcharts, and mind maps. Use analogies and visual metaphors."
    },
    "exam_focused": {
        "name": "Exam Preparation",
        "description": "Tailored for test preparation",
        "prompt_modifier": "Focus on exam-relevant content, important formulas, frequently asked questions, and exam tips. Highlight what's most likely to appear in tests."
    }
}

# --- Helper Functions ---

def extract_text_with_structure(file_stream: io.BytesIO) -> dict:
    """Extract text from PPT with slide structure preserved"""
    try:
        presentation = Presentation(file_stream)
        slides_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = {
                "slide_number": slide_num,
                "title": "",
                "content": []
            }
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                # Try to identify title
                if shape == slide.shapes.title or (hasattr(shape, 'text') and slide_num == 1):
                    slide_text["title"] = shape.text
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text["content"].append(paragraph.text.strip())
            
            if slide_text["title"] or slide_text["content"]:
                slides_content.append(slide_text)
        
        return {
            "total_slides": len(presentation.slides),
            "slides": slides_content,
            "full_text": "\n".join([
                f"Slide {s['slide_number']}: {s['title']}\n" + "\n".join(s['content']) 
                for s in slides_content
            ])
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process PowerPoint: {str(e)}")

async def generate_targeted_study_material(
    ppt_content: dict, 
    weak_spots: List[str], 
    material_type: str,
    difficulty_level: str,
    additional_instructions: str = ""
) -> str:
    """Generate study material with focus on weak spots"""
    
    if not GOOGLE_API_KEY:
        raise HTTPException(status_code=500, detail="Gemini API key not configured")
    
    model = genai.GenerativeModel(model_name='gemini-2.5-flash-preview-05-20')
    
    # Get material type configuration
    material_config = MATERIAL_TYPES.get(material_type, MATERIAL_TYPES["comprehensive"])
    
    # Build weak spots emphasis
    weak_spots_text = ""
    if weak_spots and any(weak_spots):
        weak_spots_text = f"""
        CRITICAL STUDENT WEAK SPOTS:
        The student particularly struggles with the following concepts:
        {', '.join([f'"{spot}"' for spot in weak_spots if spot.strip()])}
        
        IMPORTANT INSTRUCTIONS FOR WEAK SPOTS:
        1. Provide EXTRA detailed explanations for these topics
        2. Include multiple examples for each weak spot
        3. Break down complex concepts into simpler steps
        4. Explain prerequisites needed to understand these topics
        5. Include common misconceptions and clarifications
        6. Add practice questions specifically targeting these areas
        7. Use analogies and real-world applications
        """
    
    # Difficulty level instructions
    difficulty_instructions = {
        "beginner": "Use simple language, avoid jargon, explain everything from basics.",
        "intermediate": "Balance between foundational concepts and advanced topics.",
        "advanced": "Include complex concepts, theoretical depth, and advanced applications."
    }
    
    prompt = f"""
        You are an expert educational content creator. Generate high-quality study material based on the PowerPoint content provided.

        MATERIAL TYPE: {material_config['name']}
        STYLE INSTRUCTION: {material_config['prompt_modifier']}

        DIFFICULTY LEVEL: {difficulty_level}
        {difficulty_instructions.get(difficulty_level, "")}

        {weak_spots_text}

        ADDITIONAL TUTOR INSTRUCTIONS:
        {additional_instructions if additional_instructions else "None provided"}

        POWERPOINT CONTENT:
        Total Slides: {ppt_content['total_slides']}
        Content:
        {ppt_content['full_text']}

        FORMATTING REQUIREMENTS:
        1. Use clear hierarchical structure with headers
        2. Include relevant examples and explanations
        3. For weak spots, provide extra detail and clarity
        4. Use formatting markers: **bold** for emphasis, *italic* for terms
        5. Number important points and use bullet points for lists
        6. Include a summary section at the end
        7. Add "Key Takeaways" for each major topic
        8. If relevant, include formulas, definitions, and important dates

        OUTPUT LENGTH RESTRICTIONS:
        - The response must fit within a maximum of 2 pages in a standard PDF (≈ 600–800 words total).
        - Be concise and avoid unnecessary repetition.
        - Prioritize clarity and weak spot coverage over breadth.
        - Do not expand on all content; only cover {weak_spots_text} as specified.

        Remember: focus only on the weak spots, and compress explanations so the entire material fits within 2 PDF pages.
        """

    
    try:
        response = await model.generate_content_async(prompt)
        return response.text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate study material: {str(e)}")

def create_professional_pdf(content: str, metadata: dict, output_path: str):
    """Create a well-formatted PDF with proper structure"""
    
    import re
    from xml.sax.saxutils import escape
    
    class NumberedCanvas(canvas.Canvas):
        def __init__(self, *args, **kwargs):
            canvas.Canvas.__init__(self, *args, **kwargs)
            self._saved_page_states = []

        def showPage(self):
            self._saved_page_states.append(dict(self.__dict__))
            self._startPage()

        def save(self):
            num_pages = len(self._saved_page_states)
            for state in self._saved_page_states:
                self.__dict__.update(state)
                self.draw_page_number(num_pages)
                canvas.Canvas.showPage(self)
            canvas.Canvas.save(self)

        def draw_page_number(self, page_count):
            self.setFont("Helvetica", 9)
            self.setFillColor(grey)
            self.drawRightString(
                letter[0] - inch * 0.75,
                inch * 0.75,
                f"Page {self._pageNumber} of {page_count}"
            )
            # Add header
            self.drawString(inch * 0.75, letter[1] - inch * 0.5, 
                          f"Study Material - {metadata.get('material_type', 'Notes')}")
    
    def convert_markdown_to_html(text):
        """Properly convert markdown to HTML tags"""
        # First escape any existing XML/HTML special characters
        text = escape(text)
        
        # Convert bold markdown (**text**) to HTML
        text = re.sub(r'\*\*([^*]+)\*\*', r'<b>\1</b>', text)
        
        # Convert italic markdown (*text*) to HTML
        # Use negative lookbehind and lookahead to avoid matching bold markers
        text = re.sub(r'(?<!\*)\*([^*]+)\*(?!\*)', r'<i>\1</i>', text)
        
        return text
    
    # Create the PDF document
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )
    
    # Create custom styles
    styles = getSampleStyleSheet()
    custom_styles = {
        'CustomTitle': ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            textColor=HexColor('#1a237e'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ),
        'Metadata': ParagraphStyle(
            'Metadata',
            parent=styles['Normal'],
            fontSize=10,
            textColor=grey,
            alignment=TA_CENTER,
            spaceAfter=20
        ),
        'SectionHeader': ParagraphStyle(
            'SectionHeader',
            parent=styles['Heading1'],
            fontSize=18,
            textColor=HexColor('#1976d2'),
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold'
        ),
        'SubHeader': ParagraphStyle(
            'SubHeader',
            parent=styles['Heading2'],
            fontSize=14,
            textColor=HexColor('#424242'),
            spaceAfter=10,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        ),
        'NormalText': ParagraphStyle(
            'NormalText',
            parent=styles['Normal'],
            fontSize=11,
            textColor=black,
            alignment=TA_JUSTIFY,
            spaceAfter=8,
            leading=14
        ),
        'BulletText': ParagraphStyle(
            'BulletText',
            parent=styles['Normal'],
            fontSize=11,
            leftIndent=20,
            spaceAfter=6,
            leading=14
        ),
        'ImportantNote': ParagraphStyle(
            'ImportantNote',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#c62828'),
            borderColor=HexColor('#ffebee'),
            borderWidth=1,
            borderPadding=10,
            backColor=HexColor('#ffebee'),
            spaceAfter=12
        ),
        'KeyPoint': ParagraphStyle(
            'KeyPoint',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#1b5e20'),
            leftIndent=10,
            borderLeftColor=HexColor('#4caf50'),
            borderLeftWidth=3,
            spaceAfter=8
        )
    }
    
    story = []
    
    # Add title and metadata
    story.append(Paragraph("Personalized Study Material", custom_styles['CustomTitle']))
    
    # Add metadata
    meta_text = f"Generated on {metadata.get('date', datetime.now().strftime('%B %d, %Y'))}<br/>"
    meta_text += f"Material Type: {metadata.get('material_type', 'Study Notes')}<br/>"
    if metadata.get('weak_spots'):
        meta_text += f"Focus Areas: {', '.join(metadata['weak_spots'][:3])}"
    story.append(Paragraph(meta_text, custom_styles['Metadata']))
    story.append(Spacer(1, 20))
    
    # Process content
    lines = content.split('\n')
    in_list = False
    
    for line in lines:
        line = line.strip()
        if not line:
            if in_list:
                story.append(Spacer(1, 6))
                in_list = False
            continue
        
        # Headers
        if line.startswith('###'):
            header_text = convert_markdown_to_html(line.replace('#', '').strip())
            story.append(Paragraph(header_text, custom_styles['SubHeader']))
        elif line.startswith('##'):
            header_text = convert_markdown_to_html(line.replace('#', '').strip())
            story.append(Paragraph(header_text, custom_styles['SectionHeader']))
        elif line.startswith('#'):
            header_text = convert_markdown_to_html(line.replace('#', '').strip())
            story.append(Paragraph(header_text, custom_styles['CustomTitle']))
        
        # Key takeaways or important notes
        elif 'key takeaway' in line.lower() or 'important:' in line.lower():
            formatted_line = convert_markdown_to_html(line)
            story.append(Paragraph(formatted_line, custom_styles['KeyPoint']))
        
        # Bullet points (but not italic markers)
        elif line.startswith(('•', '-', '→')) or (line.startswith('*') and len(line) > 1 and line[1] == ' '):
            in_list = True
            if line.startswith('*'):
                bullet_text = convert_markdown_to_html(line[2:].strip())
            else:
                bullet_text = convert_markdown_to_html(line[1:].strip())
            story.append(Paragraph(f"• {bullet_text}", custom_styles['BulletText']))
        
        # Numbered lists
        elif len(line) > 2 and line[0].isdigit() and (line[1] == '.' or (len(line) > 3 and line[1:3] == '. ')):
            in_list = True
            formatted_line = convert_markdown_to_html(line)
            story.append(Paragraph(formatted_line, custom_styles['BulletText']))
        
        # Normal text
        else:
            # Convert markdown to HTML
            formatted_line = convert_markdown_to_html(line)
            try:
                story.append(Paragraph(formatted_line, custom_styles['NormalText']))
            except Exception as e:
                # If there's still an error, clean the text more aggressively
                clean_line = escape(line.replace('*', '').replace('#', ''))
                story.append(Paragraph(clean_line, custom_styles['NormalText']))
    
    # Build PDF with custom canvas for page numbers
    try:
        doc.build(story, canvasmaker=NumberedCanvas)
    except Exception as e:
        # If PDF generation fails, create a simple text version
        print(f"Warning: Complex formatting failed, creating simple PDF: {e}")
        story = []
        story.append(Paragraph("Personalized Study Material", custom_styles['CustomTitle']))
        story.append(Spacer(1, 20))
        
        # Add content as plain text
        clean_content = re.sub(r'[*#]', '', content)
        for para in clean_content.split('\n\n'):
            if para.strip():
                try:
                    story.append(Paragraph(escape(para.strip()), custom_styles['NormalText']))
                    story.append(Spacer(1, 12))
                except:
                    pass
        
        doc.build(story)

# --- API Endpoints ---

@app.post("/generate-study-material/")
async def generate_study_material_endpoint(
    file: UploadFile = File(...),
    weak_spots: str = Form("", description="Comma-separated list of topics the student struggles with"),
    material_type: str = Form("comprehensive", description="Type of study material"),
    difficulty_level: str = Form("intermediate", description="Difficulty level of content"),
    additional_instructions: str = Form("", description="Additional instructions from tutor")
):
    """
    Generate personalized study material based on student weak spots
    """
    
    # Validate file type
    if not file.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="Please upload a .pptx file")
    
    # Read and process file
    file_contents = await file.read()
    file_stream = io.BytesIO(file_contents)
    
    # Extract content with structure
    ppt_content = extract_text_with_structure(file_stream)
    
    if not ppt_content['full_text'].strip():
        raise HTTPException(status_code=400, detail="PowerPoint appears to be empty")
    
    # Process weak spots
    weak_spots_list = [spot.strip() for spot in weak_spots.split(',') if spot.strip()] if weak_spots else []
    
    # Generate study material
    try:
        generated_content = await generate_targeted_study_material(
            ppt_content=ppt_content,
            weak_spots=weak_spots_list,
            material_type=material_type,
            difficulty_level=difficulty_level,
            additional_instructions=additional_instructions
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Generation failed: {str(e)}")
    
    # Create PDF
    output_filename = f"study_material_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    
    metadata = {
        'date': datetime.now().strftime('%B %d, %Y'),
        'material_type': MATERIAL_TYPES.get(material_type, {}).get('name', 'Study Notes'),
        'weak_spots': weak_spots_list,
        'difficulty': difficulty_level
    }
    
    try:
        create_professional_pdf(generated_content, metadata, output_filename)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF creation failed: {str(e)}")
    
    # Return file
    return FileResponse(
        path=output_filename,
        media_type='application/pdf',
        filename=output_filename,
        headers={
            "Content-Disposition": f"attachment; filename={output_filename}"
        }
    )

@app.get("/api/material-types")
async def get_material_types():
    """Get available material types and their descriptions"""
    return JSONResponse(content=MATERIAL_TYPES)

@app.post("/api/analyze-ppt")
async def analyze_ppt(file: UploadFile = File(...)):
    """Analyze PPT and suggest potential weak spots based on content"""
    if not file.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="Please upload a .pptx file")
    
    file_contents = await file.read()
    file_stream = io.BytesIO(file_contents)
    ppt_content = extract_text_with_structure(file_stream)
    
    # Extract key topics from slides
    topics = []
    for slide in ppt_content['slides']:
        if slide['title']:
            topics.append(slide['title'])
    
    return JSONResponse(content={
        "total_slides": ppt_content['total_slides'],
        "main_topics": topics[:10],  # Return top 10 topics
        "content_preview": ppt_content['full_text'][:500]  # First 500 chars
    })

@app.get("/", response_class=HTMLResponse)
async def read_root():
    """Serve the main HTML page"""
    try:
        with open("index.html", "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    except FileNotFoundError:
        # Return a default HTML if file not found
        return HTMLResponse(
            content="<h1>Error</h1><p>index.html not found. Please ensure the file exists.</p>",
            status_code=404
        )

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "api_key_configured": bool(GOOGLE_API_KEY)}